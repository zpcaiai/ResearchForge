"""What the output plane refuses to draw, present or ship.

These tests are written against the refusals. A figure generator, a deck builder
and an exporter all have trivially satisfiable happy paths — they emit files. The
only thing worth testing is whether they emit files they should not: a figure
whose numbers disagree with the analysis, a slide whose number came from nowhere,
a release that ships an artifact with no lineage or a placeholder produced by the
offline stub.
"""
import io
import json
import statistics
import xml.etree.ElementTree as ET
from pathlib import Path

import pytest
from jsonschema import Draft202012Validator

from researchforge import skills as _skills  # noqa: F401  registers implementations
from researchforge.artifacts import ArtifactStore
from researchforge.errors import GateBlocked
from researchforge.providers import OfflineStubProvider, QuotaLedger
from researchforge.provenance import ProvenanceLog, sha256_file
from researchforge.skill import Context, get

SCHEMAS = Path(__file__).resolve().parents[2] / "schemas"

DISCLOSURE = ("The comparison baseline was not reproduced locally (level RL1). All comparisons "
              "are against numbers as reported by the original authors, under a different "
              "environment. Observed shortfalls: HARDWARE_UNAVAILABLE.")

VALUES = {"main": [0.81, 0.8213, 0.83], "probe": [0.77, 0.79, 0.7801]}


# --------------------------------------------------------------------------
# harness
# --------------------------------------------------------------------------
def make_ctx(tmp_path, **config):
    prov = ProvenanceLog(tmp_path)
    store = ArtifactStore(tmp_path, SCHEMAS, run_id="test", provenance=prov)
    return Context(project=tmp_path, run_id="test", mode="auto", store=store, prov=prov,
                   quota=QuotaLedger(tmp_path / "quota.jsonl"), model=OfflineStubProvider(),
                   scholarly=[], config=config, offline=True)


def comparison_mode(mode="CM_REPORTED", level="RL1", *, disclosure=True):
    forbidden = {"CM_NONE": ["outperforms", "improves over", "state-of-the-art"],
                 "CM_REPORTED": ["we reproduce", "our measured baseline",
                                 "under identical conditions"]}.get(mode, [])
    return {
        "mode": mode, "derived_from_level": level,
        "admissible_idea_modes": ["explain_diagnose", "benchmark_evaluate"],
        "forbidden_claim_patterns": forbidden,
        "disclosure_required": {"required": disclosure,
                                "text_template": DISCLOSURE if disclosure else "",
                                "must_appear_in": ["experimental setup", "limitations"]
                                if disclosure else []},
        "substitute_baseline": None, "approved_by": None, "autonomous_decision_id": None,
    }


def ledger():
    rows = []
    for branch, vals in VALUES.items():
        for seed, v in enumerate(vals):
            rows.append({"run_id": f"{branch}-E-1-s{seed}", "experiment_id": "E-1", "status": "ok",
                         "metrics": {"accuracy": v}, "artifacts": [],
                         "provenance": {"seed": seed, "branch": branch,
                                        "evaluator_version": "evaluate.py@1.0.0",
                                        "environment_digest": "sha256:env-a"}})
    return rows


def _group(branch, metric="accuracy", claim="C-1"):
    vals = VALUES[branch]
    mean = statistics.fmean(vals)
    sd = statistics.stdev(vals)
    half = 4.302652729911275 * sd / (len(vals) ** 0.5)          # t(0.975, df=2)
    return {"group_id": f"{branch}::{metric}", "branch": branch, "metric": metric,
            "kind": "confirmatory", "seeds": [0, 1, 2],
            "run_ids": [f"{branch}-E-1-s{i}" for i in range(len(vals))],
            "values": vals, "n": len(vals), "mean": mean, "median": statistics.median(vals),
            "min": min(vals), "max": max(vals), "sd": sd, "sem": sd / (len(vals) ** 0.5),
            "ci95": [mean - half, mean + half],
            "strata": [{"evaluator_version": "evaluate.py@1.0.0",
                        "environment_digest": "sha256:env-a"}]}


def analysis_results(claim="C-1"):
    groups = [_group(b) for b in VALUES]
    return {
        "run_id": "test", "generated_at": 0.0, "analysis_question": "does the probe isolate it",
        "ledger_digest": "sha256:ledger", "n_runs": 6, "n_successful": 6, "n_failed": 0,
        "failure_rate": 0.0, "confirmatory_metrics": ["accuracy"], "exploratory_metrics": [],
        "groups": groups,
        "reported": [{"reported_id": f"r::{g['group_id']}::mean", "group_id": g["group_id"],
                      "branch": g["branch"], "metric": g["metric"], "statistic": "mean",
                      "value": g["mean"], "interval": g["ci95"], "claim_ids": [claim]}
                     for g in groups],
        "comparisons": [],
    }


GRAPH = [{"claim_id": "C-1", "claim_text": "The probe isolates the failure mode.",
          "claim_type": "empirical", "status": "SUPPORTED", "conflicts": [],
          "support_edges": [{"result_id": "E-1", "relation": "supports"}]}]


def spine(statement="The probe isolates the failure mode."):
    return {"thesis": "A diagnostic account of the failure mode",
            "target_venue": "unspecified venue", "comparison_mode": "CM_REPORTED",
            "sections": ["Results", "Limitations"],
            "claims": [{"claim_id": "C-1", "statement": statement, "kind": "empirical",
                        "section": "Results",
                        "evidence": {"ref_ids": [], "result_ids": ["E-1"],
                                     "graph_status": "SUPPORTED"}}],
            "argument_chain": ["C-1"]}


def figure_plan(caption="Per-seed accuracy for each branch.", claim_id="C-1"):
    return {"figures": [{"figure_id": "F-01", "claim_id": claim_id, "kind": "results",
                         "caption": caption, "data_source": {"result_ids": ["E-1"]},
                         "status": "PLANNED"}],
            "comparison_mode": "CM_REPORTED", "no_figure_without_a_claim": True}


def seed_figures(ctx, *, plan=None, results=None):
    s = ctx.store
    s.write("claim-evidence-graph", "evidence_graph", GRAPH)
    s.write("data-analyst", "analysis_results", analysis_results() if results is None else results)
    s.write("data-analyst", "analysis_plots", {"rendered": True, "plots": []})
    s.write("manuscript-builder", "manuscript_spine", spine())
    s.write("manuscript-builder", "figure_plan", figure_plan() if plan is None else plan)
    return ctx


def seed_deck(ctx, *, statement="The probe isolates the failure mode.", mode="CM_REPORTED",
              level="RL1"):
    s = ctx.store
    s.write("manuscript-builder", "manuscript_spine", spine(statement))
    s.write("manuscript-builder", "draft_manifest",
            {"sections": ["Results"], "paragraphs": [], "claim_index": {},
             "every_paragraph_bound_to_a_claim": True, "_synthetic": False})
    s.write("manuscript-builder", "manuscript_draft", "\\section{Results}\n" + DISCLOSURE + "\n")
    s.write("experiment-runner", "experiment_ledger", ledger())
    s.write("integrity-auditor", "meta_analysis", {"id": "MA-1", "evidence_locked": True})
    s.write("reproduction-fallback-planner", "comparison_mode", comparison_mode(mode, level))
    return ctx


# --------------------------------------------------------------------------
# figure-factory
# --------------------------------------------------------------------------
def test_figures_are_vector_editable_and_bound_to_a_claim(tmp_path):
    ctx = seed_figures(make_ctx(tmp_path))
    get("figure-factory")(ctx)

    svgs = sorted((tmp_path / "figures" / "selected").glob("*.svg"))
    assert svgs, "figure-factory produced no SVG"
    root = ET.fromstring(svgs[0].read_text())
    ns = "{http://www.w3.org/2000/svg}"
    # Never flatten: a raster figure cannot be corrected, only redrawn.
    assert not list(root.iter(f"{ns}image"))
    # Text stays text, or the labels are an outline nobody can edit or search.
    assert list(root.iter(f"{ns}text"))

    emap = json.loads((tmp_path / "figures" / "element_map.json").read_text())
    fig = emap["figures"]["F-01"]
    assert fig["claim_id"] == "C-1"
    assert all(e["binds_to"] for e in fig["elements"])
    means = [e for e in fig["elements"] if e["role"] == "mean"]
    plotted = sorted(round(e["values"][0], 6) for e in means)
    expected = sorted(round(statistics.fmean(v), 6) for v in VALUES.values())
    assert plotted == expected

    trace = json.loads((tmp_path / "figures" / "generation_trace.json").read_text())
    assert trace["figures"][0]["rasterized"] is False
    # The editable copy is the shipped copy; two files would be two truths.
    assert (tmp_path / "figures" / "editable" / "F-01.svg").read_bytes() == svgs[0].read_bytes()


def test_figure_whose_numbers_disagree_with_the_analysis_is_refused(tmp_path):
    ctx = seed_figures(make_ctx(tmp_path),
                       plan=figure_plan(caption="Accuracy reaches 0.99 on the probed branch."))
    with pytest.raises(GateBlocked) as e:
        get("figure-factory")(ctx)
    assert e.value.gate == "figure_data_integrity"
    assert "0.99" in str(e.value)
    # Refusal means nothing was drawn, not that a warning was attached to a drawing.
    assert not ctx.store.exists("selected_figure")
    assert not ctx.store.exists("editable_svg")


def test_figure_that_names_no_known_claim_is_refused(tmp_path):
    ctx = seed_figures(make_ctx(tmp_path), plan=figure_plan(claim_id="C-99"))
    with pytest.raises(GateBlocked) as e:
        get("figure-factory")(ctx)
    assert "C-99" in str(e.value)
    assert not ctx.store.exists("selected_figure")


def test_figure_factory_refuses_without_the_analysis(tmp_path):
    ctx = make_ctx(tmp_path)
    ctx.store.write("manuscript-builder", "manuscript_spine", spine())
    ctx.store.write("manuscript-builder", "figure_plan", figure_plan())
    with pytest.raises(GateBlocked) as e:
        get("figure-factory")(ctx)
    assert "analysis_results" in str(e.value)


# --------------------------------------------------------------------------
# deck-factory
# --------------------------------------------------------------------------
def build_deck(ctx):
    get("figure-factory")(ctx)
    get("deck-factory")(ctx)
    from pptx import Presentation
    return Presentation(io.BytesIO((ctx.project / "slides" / "defense.pptx").read_bytes()))


def test_deck_is_native_powerpoint_not_a_stack_of_images(tmp_path):
    ctx = seed_deck(seed_figures(make_ctx(tmp_path)))
    prs = build_deck(ctx)

    from pptx.enum.shapes import MSO_SHAPE_TYPE

    text_frames, tables, pictures = 0, 0, 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_frames += 1
            if getattr(shape, "has_table", False):
                tables += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
    assert len(prs.slides) >= 5
    assert text_frames >= len(prs.slides), "every slide must carry real, editable text"
    assert tables >= 1, "the results table must be a real table, not a picture of one"
    # The spec forbids a deck of full-slide images; so does auditability.
    assert pictures == 0
    assert all(s.has_notes_slide for s in prs.slides)

    manifest = json.loads((tmp_path / "slides" / "deck_manifest.json").read_text())
    assert manifest["full_slide_images"] == 0
    assert manifest["every_quantitative_element_bound"] is True


def test_every_number_on_a_slide_names_the_artifact_it_came_from(tmp_path):
    ctx = seed_deck(seed_figures(make_ctx(tmp_path)))
    build_deck(ctx)
    evidence = json.loads((tmp_path / "slides" / "slide_evidence.json").read_text())
    quantitative = [e for e in evidence["elements"] if e["quantitative"]]
    assert quantitative
    assert all(e["bound_to"] for e in quantitative)
    assert "experiment_ledger" in evidence["artifacts_referenced"]
    assert evidence["unbound"] == []


def test_slide_number_bound_to_no_artifact_blocks_the_deck(tmp_path):
    ctx = seed_figures(make_ctx(tmp_path))
    seed_deck(ctx, statement="The probe raises accuracy to 0.99 in every setting.")
    with pytest.raises(GateBlocked) as e:
        build_deck(ctx)
    assert e.value.gate == "slide_evidence_binding"
    assert "0.99" in str(e.value)
    assert not ctx.store.exists("defense_deck")


def test_comparative_language_is_removed_under_cm_none(tmp_path):
    ctx = seed_figures(make_ctx(tmp_path))
    seed_deck(ctx, statement="The probe outperforms the published baseline.",
              mode="CM_NONE", level="RL0")
    build_deck(ctx)
    manifest = json.loads((tmp_path / "slides" / "deck_manifest.json").read_text())
    assert manifest["removed_for_comparison_mode"]
    outline = (tmp_path / "slides" / "slide_outline.md").read_text()
    assert "outperforms" not in outline


def test_deck_refuses_without_figures(tmp_path):
    ctx = seed_deck(make_ctx(tmp_path))
    with pytest.raises(GateBlocked) as e:
        get("deck-factory")(ctx)
    assert "selected_figure" in str(e.value)


# --------------------------------------------------------------------------
# release-gate
# --------------------------------------------------------------------------
CLEAN_GATE = {"verdict": "PASS", "submission_permitted": True, "decided_at": 0.0,
              "run_id": "test", "comparison_mode": "CM_REPORTED", "counts": {}, "checks": [],
              "blockers": [], "model_consulted": False, "synthetic_inputs": False,
              "human_reviewed": False}


def seed_release(ctx, *, gate=None, review=None, stats=None, mode="CM_REPORTED", level="RL1",
                 draft=None, manifest_extra=()):
    s = ctx.store
    s.write("claim-citation-auditor", "integrity_gate", CLEAN_GATE if gate is None else gate)
    s.write("claim-citation-auditor", "submission_blockers",
            "# Submission blockers — verdict: PASS\n\nNo blocker survived the audit.\n")
    s.write("claim-citation-auditor", "claim_audit",
            [{"claim_id": "C-1", "locator": "Results#p0", "verdict": "SUPPORTED",
              "text": "The probe isolates the failure mode.", "citations": []}])
    s.write("claim-citation-auditor", "citation_audit",
            "# Citation audit\n\nNo citation in this draft is unresolved.\n")
    s.write("review-simulator", "review_report",
            {"recommendation": "REVISE_THEN_SUBMIT",
             "recommendation_rationale": "no blocker survived the audit",
             "simulated_reviews": [], "grounded_points": [], "_synthetic": False}
            if review is None else review)
    s.write("integrity-auditor", "stats_audit",
            {"findings": [], "severity_counts": {},
             "evidence_lock": {"blocked": False, "blocked_by": [], "blocked_claims": []}}
            if stats is None else stats)
    s.write("reproduction-fallback-planner", "comparison_mode", comparison_mode(mode, level))
    s.write("result-reproducer", "source_repro_report",
            {"target_paper_id": "p-1", "target_kind": "source_paper", "level": level,
             "assessed_at_run_id": "test", "claim_comparisons": [], "failure_codes": [],
             "environment_digest": "not-captured", "timebox_seconds": 60.0,
             "timebox_exhausted": False, "human_reviewed": False})
    s.write("manuscript-builder", "manuscript_draft",
            "\\section{Results}\nThe probe isolates the failure mode.\n"
            "\\section{Limitations}\n" + DISCLOSURE + "\n" if draft is None else draft)
    s.write("citation-resolver", "bibliography", "@misc{R-001,\n  note = {A real work},\n}\n")
    s.write("finding-memory", "findings",
            [{"finding_id": "F-001", "statement": "The probe isolates the failure mode.",
              "result_ids": ["E-1"]}])
    s.write("finding-memory", "finding_memory_graph", {"nodes": ["F-001"], "edges": []})
    s.write("experiment-runner", "experiment_ledger", ledger())
    s.write("deck-factory", "defense_deck", b"PK\x03\x04 stand-in bytes for a built deck")
    digest = sha256_file(ctx.project / "experiment_ledger.jsonl")
    s.write("experiment-runner", "artifact_manifest",
            {"project_id": ctx.project.name,
             "artifacts": [{"artifact_id": "experiment_ledger", "path": "experiment_ledger.jsonl",
                            "sha256": digest, "kind": "ledger", "parents": []},
                           *manifest_extra]})
    return ctx


def release(ctx):
    get("release-gate")(ctx)
    return (json.loads((ctx.project / "release" / "release_manifest.json").read_text()),
            json.loads((ctx.project / "release" / "_manifest.json").read_text()),
            (ctx.project / "release" / "release_report.md").read_text())


def kinds(bundle):
    return {b["kind"] for b in bundle["blockers"]}


def test_clean_release_is_an_assisted_draft_and_never_submission_ready(tmp_path):
    ctx = seed_release(make_ctx(tmp_path))
    manifest, bundle, report = release(ctx)

    assert bundle["released"] is True
    assert bundle["release_status"] == "ASSISTED_DRAFT"
    # A system that writes manuscripts must not be the thing that declares them ready.
    assert manifest["submission_ready"] is False
    assert manifest["ai_participation"]["human_verified"] is False
    assert "ResearchForge" in manifest["ai_participation"]["statement"]
    assert "assisted draft" in report.lower()
    # the deliverables actually landed in the bundle
    assert (tmp_path / "release" / "paper" / "main.tex").exists()
    assert any(d["artifact_id"] == "manuscript_draft" for d in bundle["deliverables"])
    schema = json.loads((SCHEMAS / "ArtifactManifest.schema.json").read_text())
    Draft202012Validator(schema).validate(manifest)


def test_unresolved_blocker_stops_the_release(tmp_path):
    gate = {**CLEAN_GATE, "verdict": "BLOCK", "submission_permitted": False,
            "blockers": [{"blocker_id": "B-001", "kind": "claim_fabricated", "claim_id": "C-1",
                          "locator": "Results#p0", "detail": "the number 0.99 matches no run",
                          "remediation": "delete the number"}]}
    ctx = seed_release(make_ctx(tmp_path), gate=gate)
    manifest, bundle, report = release(ctx)

    assert bundle["released"] is False
    assert bundle["release_status"] == "BLOCKED"
    assert "unresolved_claim_fabricated" in kinds(bundle)
    assert bundle["deliverables"] == []
    assert not (tmp_path / "release" / "paper" / "main.tex").exists()
    assert "refused" in report.lower()
    assert manifest["released"] is False


def test_synthetic_artifact_stops_the_release(tmp_path):
    ctx = seed_release(make_ctx(tmp_path),
                       review={"recommendation": "REVISE_THEN_SUBMIT",
                               "recommendation_rationale": "clean", "simulated_reviews": [],
                               "_synthetic": True})
    _manifest, bundle, _report = release(ctx)
    assert bundle["released"] is False
    assert "synthetic_artifact_in_release" in kinds(bundle)


def test_artifact_without_provenance_stops_the_release(tmp_path):
    ctx = make_ctx(tmp_path)
    rogue = tmp_path / "analysis" / "rogue_table.csv"
    rogue.parent.mkdir(parents=True, exist_ok=True)
    rogue.write_text("metric,value\naccuracy,0.99\n")
    seed_release(ctx, manifest_extra=[{"artifact_id": "rogue_table",
                                       "path": "analysis/rogue_table.csv",
                                       "sha256": sha256_file(rogue), "kind": "table",
                                       "parents": []}])
    _manifest, bundle, _report = release(ctx)
    assert bundle["released"] is False
    assert "artifact_without_provenance" in kinds(bundle)
    assert any("rogue_table" in b["detail"] for b in bundle["blockers"])


def test_artifact_edited_after_it_was_audited_stops_the_release(tmp_path):
    ctx = seed_release(make_ctx(tmp_path))
    # What ships must be what was audited; a later edit breaks that chain silently.
    (tmp_path / "paper" / "main.tex").write_text("\\section{Results}\nAccuracy is 0.99.\n")
    _manifest, bundle, _report = release(ctx)
    assert bundle["released"] is False
    assert "artifact_digest_drift" in kinds(bundle)


def test_missing_disclosure_under_cm_reported_stops_the_release(tmp_path):
    ctx = seed_release(make_ctx(tmp_path),
                       draft="\\section{Results}\nThe probe isolates the failure mode.\n")
    _manifest, bundle, _report = release(ctx)
    assert bundle["released"] is False
    assert "missing_disclosure" in kinds(bundle)


def test_comparative_claim_does_not_survive_cm_none(tmp_path):
    ctx = seed_release(
        make_ctx(tmp_path), mode="CM_NONE", level="RL0",
        draft="\\section{Results}\nOur method outperforms the published baseline.\n")
    _manifest, bundle, _report = release(ctx)
    assert bundle["released"] is False
    assert "comparative_claim_under_CM_NONE" in kinds(bundle)


def test_release_gate_refuses_when_an_upstream_audit_is_missing(tmp_path):
    ctx = seed_release(make_ctx(tmp_path))
    (tmp_path / "review" / "review_report.json").unlink()
    with pytest.raises(GateBlocked) as e:
        get("release-gate")(ctx)
    assert e.value.gate == "release_inputs"
    assert "review_report" in str(e.value)


def test_upstream_statistics_blocker_is_carried_into_the_release(tmp_path):
    ctx = seed_release(make_ctx(tmp_path),
                       stats={"findings": [{"finding_id": "S-1", "severity": "BLOCKER",
                                            "code": "UNDERPOWERED",
                                            "message": "n=3 cannot support this contrast",
                                            "remediation": "run more seeds"}],
                              "evidence_lock": {"blocked": True, "blocked_by": ["S-1"],
                                                "blocked_claims": ["C-1"]}})
    _manifest, bundle, _report = release(ctx)
    assert bundle["released"] is False
    assert {"stats_audit_unresolved", "evidence_not_locked"} <= kinds(bundle)
