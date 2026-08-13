"""Output plane: figures, the defense deck, and the release gate.

Three skills, and the third is the reason the first two are constrained the way
they are.

`figure-factory` draws real vector figures from the audited analysis results and
from nothing else. Every figure names the claim it supports, and every number it
draws is read out of `analysis_results` at draw time. A figure whose numbers
disagree with the analysis is a refusal, not a line in a validation report: a
figure is the artifact a reader believes fastest and checks last, so a wrong one
is worse than a missing one. Nothing is flattened to raster while a vector path
exists — a raster figure cannot be corrected, only redrawn, and no reviewer can
read a number out of a picture of one.

`deck-factory` builds a native PowerPoint: real text frames, real shapes, real
tables. A deck of full-slide images is forbidden by the spec, and would be
forbidden anyway — it is unauditable, uneditable and unreadable to a screen
reader. Every quantitative element on every slide is bound to a project artifact
id in `slide_evidence`, and a number that binds to nothing stops the build,
because a defense slide is where an unbacked number is stated out loud to people
who cannot check it in the room.

`release-gate` is the last place a fabricated number can be caught. It reads
every upstream audit and refuses to release while any blocker is unresolved; it
refuses to release an artifact that has no provenance record or no digest, since
an artifact with no lineage cannot support a claim; it treats anything the
offline stub produced as a hard blocker; and it defaults the release to an
assisted draft carrying an AI-participation disclosure rather than to something
marked submission-ready. The machine wrote most of this. The manifest says so.
"""
from __future__ import annotations

import hashlib
import io
import json
import re
import shutil
import time
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any

from ..errors import GateBlocked, SchemaViolation
from ..generated import ARTIFACTS
from ..provenance import sha256_file
from ..skill import Context, Skill, SkillResult, register
from .writing import COMPARATIVE_RE, _ledger_index, _match_number, _norm, _quantities, _walk_numbers

SVG_NS = "http://www.w3.org/2000/svg"

#: Keys that mark a payload as produced by the offline stub provider. Any of them
#: true anywhere in the release chain is a hard blocker: stub output is a
#: structural placeholder, and a placeholder that reaches a reader is a fabrication
#: no matter how carefully it was labelled upstream.
SYNTHETIC_KEYS = ("synthetic", "_synthetic", "is_synthetic")


# --------------------------------------------------------------------------
# shared helpers
# --------------------------------------------------------------------------
def _missing(ctx: Context, ids: tuple[str, ...]) -> list[str]:
    return [i for i in ids if not ctx.store.exists(i)]


def _dir_manifest(ctx: Context, artifact_id: str, *, default: Any = None) -> Any:
    """Read a directory-valued artifact's manifest.

    `ArtifactStore.read` cannot open a directory, so the `_manifest.json` inside
    it is read directly. Existence still goes through the store, so a directory
    artifact that was never produced reports as missing rather than as empty.
    """
    if not ctx.store.exists(artifact_id):
        return default
    p = ctx.store.path_for(artifact_id) / "_manifest.json"
    return json.loads(p.read_text(encoding="utf-8"))


def _synthetic_hits(label: str, payload: Any, path: str = "") -> list[str]:
    """Every place in a payload that admits to being stub output."""
    hits: list[str] = []
    if isinstance(payload, dict):
        for k, v in payload.items():
            here = f"{path}.{k}" if path else str(k)
            if k in SYNTHETIC_KEYS and v is True:
                hits.append(f"{label}:{here}")
            hits += _synthetic_hits(label, v, here)
    elif isinstance(payload, list):
        for i, v in enumerate(payload):
            hits += _synthetic_hits(label, v, f"{path}[{i}]")
    return hits


def _csv(header: list[str], rows: list[list[Any]]) -> str:
    esc = lambda s: '"' + str(s).replace('"', '""') + '"'
    return "\n".join([",".join(header)] + [",".join(esc(c) for c in r) for r in rows]) + "\n"


def _slug(s: str) -> str:
    return re.sub(r"[^A-Za-z0-9._-]+", "-", str(s)).strip("-") or "unnamed"


def _close(a: float, b: float, tol: float = 1e-9) -> bool:
    return abs(float(a) - float(b)) <= tol


# ==========================================================================
# figure-factory
# ==========================================================================
def _value_index(results: dict[str, Any]) -> list[dict[str, Any]]:
    """Every number the audited analysis actually produced, with its origin.

    A figure may draw a number that appears in here and no other number. This is
    the same rule the claim auditor applies to prose, applied to pixels.
    """
    idx: list[dict[str, Any]] = []
    for g in results.get("groups") or []:
        base = {"group_id": g.get("group_id"), "branch": g.get("branch"),
                "metric": g.get("metric"), "kind": g.get("kind")}
        for stat in ("mean", "median", "min", "max", "sd", "sem"):
            if isinstance(g.get(stat), (int, float)):
                idx.append({**base, "statistic": stat, "value": float(g[stat])})
        for i, v in enumerate(g.get("values") or []):
            if isinstance(v, (int, float)):
                idx.append({**base, "statistic": f"values[{i}]", "value": float(v)})
        for i, v in enumerate(g.get("ci95") or []):
            if isinstance(v, (int, float)):
                idx.append({**base, "statistic": f"ci95[{i}]", "value": float(v)})
        if isinstance(g.get("n"), int):
            idx.append({**base, "statistic": "n", "value": float(g["n"])})
    for r in results.get("reported") or []:
        if isinstance(r.get("value"), (int, float)):
            idx.append({"group_id": r.get("group_id"), "branch": r.get("branch"),
                        "metric": r.get("metric"), "kind": "reported",
                        "statistic": r.get("statistic"), "value": float(r["value"]),
                        "claim_ids": r.get("claim_ids") or []})
    return idx


@register
class FigureFactory(Skill):
    """Storyboard, draw, reconstruct and validate — all against the analysis.

    Consolidates figure-storyboard, scientific-figure-generator,
    vector-figure-reconstructor and editable-svg-refiner. They were four skills
    that each held one end of the same invariant: the figure shows what the
    analysis found. Split across four skills the invariant was nobody's, so it is
    enforced here in one place, at draw time, on the numbers that are actually
    plotted.
    """

    name = "figure-factory"

    def execute(self, ctx: Context) -> SkillResult:
        gone = _missing(ctx, ("figure_plan", "analysis_results", "manuscript_spine"))
        if gone:
            raise GateBlocked(
                "figure_inputs",
                f"figure-factory cannot draw anything: {gone} have not been produced. A figure "
                f"drawn without the analysis behind it is a picture of a number nobody measured.",
                "run data-analyst and manuscript-builder first")
        plan = ctx.store.read(self.name, "figure_plan")
        results = ctx.store.read(self.name, "analysis_results")
        spine = ctx.store.read(self.name, "manuscript_spine")
        graph = ctx.store.read(self.name, "evidence_graph", default=[])
        plots = _dir_manifest(ctx, "analysis_plots", default={"rendered": False})
        venue = ctx.external("target_venue", "unspecified venue")
        style = ctx.external("figure_style_reference", {}) or {}
        warnings: list[str] = []

        index = _value_index(results)
        if not index:
            raise GateBlocked(
                "figure_inputs",
                "analysis_results carries no numbers (no groups, no reported statistics), so "
                "every figure this plan asks for would have to be invented.",
                "run data-analyst against a non-empty experiment ledger")

        known_claims = ({str(c.get("claim_id")) for c in spine.get("claims") or []} |
                        {str(e.get("claim_id")) for e in graph if e.get("claim_id")})
        storyboard, refusals = self._storyboard(plan, results, index, known_claims)
        if refusals:
            # A figure is believed faster than a sentence and checked later. One that
            # disagrees with the analysis, or that names no claim, is therefore not
            # emitted with a warning attached — nothing is emitted at all.
            raise GateBlocked(
                "figure_data_integrity",
                "figure-factory refuses to draw " + f"{len(refusals)} planned figure(s): " +
                "; ".join(f"{r['figure_id']}: {r['detail']}" for r in refusals),
                "correct paper/figure_plan.json against analysis/analysis_results.json, or "
                "re-run data-analyst so the numbers the figures need exist")

        mpl, plt = self._backend()
        trace: list[dict[str, Any]] = []
        element_map: dict[str, Any] = {}
        selected: list[dict[str, Any]] = []
        candidates_dir = ctx.store.path_for("figure_candidates")
        selected_dir = ctx.store.path_for("selected_figure")
        editable_dir = ctx.store.path_for("editable_svg")
        for d in (candidates_dir, selected_dir, editable_dir):
            d.mkdir(parents=True, exist_ok=True)

        for item in storyboard["items"]:
            variants = []
            for form in item["candidate_forms"]:
                svg, drawn = self._render(plt, item, form, style)
                drift = self._drift(drawn, item)
                if drift:
                    # The check is on what was drawn, read back off the artists, not
                    # on what the code meant to draw.
                    raise GateBlocked(
                        "figure_data_integrity",
                        f"the rendered figure {item['figure_id']} ({form}) does not agree with "
                        f"analysis_results: {drift}",
                        "this is a defect in figure-factory, not in the data; do not ship the "
                        "figure")
                name = f"{_slug(item['figure_id'])}__{form}.svg"
                (candidates_dir / name).write_bytes(svg)
                variants.append({"form": form, "file": name, "elements": drawn["elements"],
                                 "sha256": hashlib.sha256(svg).hexdigest(), "svg": svg})
            chosen = self._choose(variants, item)
            final = f"{_slug(item['figure_id'])}.svg"
            (selected_dir / final).write_bytes(chosen["svg"])
            # The selected figure and the editable figure are the same bytes. Shipping
            # a flattened "final" beside an editable "source" is how the two drift.
            (editable_dir / final).write_bytes(chosen["svg"])
            checks = self._validate_svg(chosen["svg"], chosen["elements"])
            element_map[item["figure_id"]] = {
                "file": final, "claim_id": item["claim_id"], "form": chosen["form"],
                "sha256": chosen["sha256"], "elements": chosen["elements"],
                "svg_checks": checks}
            selected.append({"figure_id": item["figure_id"], "claim_id": item["claim_id"],
                             "file": final, "form": chosen["form"], "sha256": chosen["sha256"],
                             "caption": item["caption"]})
            trace.append({
                "figure_id": item["figure_id"], "claim_id": item["claim_id"],
                "message": item["message"],
                "source": {"artifact": "analysis_results", "group_ids": item["group_ids"],
                           "run_ids": item["run_ids"]},
                "candidates": [{"form": v["form"], "file": v["file"], "sha256": v["sha256"]}
                               for v in variants],
                "selected": {"form": chosen["form"], "file": final},
                "selection_rule": item["selection_rule"],
                "renderer": f"matplotlib {mpl.__version__} (Agg, svg.fonttype=none)",
                "rasterized": False,
                "checks": checks,
            })

        recon, recon_map, recon_warning = self._reconstruct(ctx)
        if recon_warning:
            warnings.append(recon_warning)

        ctx.store.write(self.name, "figure_storyboard", storyboard)
        ctx.store.write(self.name, "figure_candidates",
                        {"figures": [{"figure_id": t["figure_id"], "candidates": t["candidates"]}
                                     for t in trace],
                         "note": "every candidate is vector; no raster variant is generated"})
        ctx.store.write(self.name, "selected_figure",
                        {"figures": selected, "target_venue": venue,
                         "format": "svg", "rasterized": False})
        ctx.store.write(self.name, "editable_svg",
                        {"figures": selected, "editable": True,
                         "text_is_text_not_paths": True,
                         "identical_to_selected": True,
                         "note": "the selected figure and the editable figure are the same "
                                 "bytes; a flattened copy beside an editable one is how the two "
                                 "silently diverge"})
        ctx.store.write(self.name, "svg_element_map",
                        {"generated_at": time.time(), "figures": element_map,
                         "binding_rule": "every element that carries a number names the "
                                         "analysis_results group and statistic it was read from"})
        ctx.store.write(self.name, "figure_generation_trace",
                        {"run_id": ctx.run_id, "generated_at": time.time(),
                         "analysis_ledger_digest": results.get("ledger_digest"),
                         "upstream_diagnostic_plots": plots.get("rendered", False),
                         "figures": trace, "reconstruction": recon_map})
        ctx.store.write(self.name, "svg_reconstruction_map", recon_map)
        ctx.store.write(self.name, "reconstructed_svg", recon)
        ctx.store.write(self.name, "figure_captions", self._captions(storyboard))
        ctx.store.write(self.name, "svg_validation_report",
                        self._validation_md(trace, element_map, results, venue))

        if not plots.get("rendered"):
            warnings.append("data-analyst rendered no diagnostic plots, so the figures below were "
                            "drawn from analysis_results alone; there is nothing to cross-read "
                            "them against visually.")
        return SkillResult(
            self.name,
            produced=["figure_storyboard", "figure_candidates", "selected_figure", "editable_svg",
                      "svg_element_map", "figure_generation_trace", "svg_reconstruction_map",
                      "reconstructed_svg", "figure_captions", "svg_validation_report"],
            warnings=warnings, next_state="WRITING",
            detail={"figures": len(selected), "candidates": sum(len(t["candidates"]) for t in trace),
                    "rasterized": False})

    # ------------------------------------------------------------------
    def _backend(self):
        try:
            import matplotlib
            matplotlib.use("Agg")
            # Text stays text. With fonttype='path' every label becomes an outline,
            # which is a raster figure's illegibility achieved in vector form: nobody
            # can edit a typo or search for an axis label.
            matplotlib.rcParams["svg.fonttype"] = "none"
            matplotlib.rcParams["svg.hashsalt"] = "researchforge"
            import matplotlib.pyplot as plt
        except Exception as e:  # noqa: BLE001 - the reason matters more than the class
            raise GateBlocked(
                "figure_renderer",
                f"matplotlib is unavailable ({e}), so no vector figure can be drawn. "
                f"figure-factory will not substitute a raster image or a placeholder.",
                "install matplotlib") from e
        return matplotlib, plt

    def _storyboard(self, plan, results, index, known_claims):
        """Message first, form second — and both bound to a claim.

        The storyboard is where a figure acquires the claim it argues. A figure
        that reaches the rendering stage without one has nothing to be right or
        wrong about.
        """
        items: list[dict[str, Any]] = []
        refusals: list[dict[str, Any]] = []
        groups = {g.get("group_id"): g for g in results.get("groups") or []}
        for i, fig in enumerate(plan.get("figures") or []):
            fid = str(fig.get("figure_id") or f"F-{i + 1:02d}")
            claim_id = fig.get("claim_id")
            if not claim_id or str(claim_id) not in known_claims:
                refusals.append({"figure_id": fid, "kind": "unbound_figure",
                                 "detail": f"names claim {claim_id!r}, which is in neither the "
                                           f"manuscript spine nor the evidence graph; a figure "
                                           f"that supports no claim cannot be checked against "
                                           f"anything"})
                continue
            gids = self._resolve_groups(fig, claim_id, results)
            if not gids:
                refusals.append({"figure_id": fid, "kind": "figure_data_missing",
                                 "detail": f"no group in analysis_results is bound to claim "
                                           f"{claim_id} or to result ids "
                                           f"{(fig.get('data_source') or {}).get('result_ids')}"})
                continue
            asserted = self._asserted_numbers(fig)
            bad = [a for a in asserted if _match_number(a, index) is None]
            if bad:
                refusals.append({"figure_id": fid, "kind": "figure_number_mismatch",
                                 "detail": f"states {[b['raw'] for b in bad]}, which match no "
                                           f"value in analysis_results ({len(index)} values). A "
                                           f"figure whose numbers disagree with the analysis is a "
                                           f"fabrication with a caption on it"})
                continue
            series = []
            for gid in gids:
                g = groups[gid]
                series.append({"group_id": gid, "branch": g.get("branch"), "metric": g.get("metric"),
                               "values": [float(v) for v in g.get("values") or []],
                               "mean": g.get("mean"), "ci95": g.get("ci95"), "n": g.get("n"),
                               "kind": g.get("kind")})
            has_interval = any(s["ci95"] for s in series)
            forms = ["per_seed_scatter"] + (["interval_plot"] if has_interval else [])
            items.append({
                "figure_id": fid, "claim_id": str(claim_id),
                "message": _norm(fig.get("caption")) or f"evidence for {claim_id}",
                "caption": _norm(fig.get("caption")) or f"Evidence for {claim_id}.",
                "metric": series[0]["metric"],
                "group_ids": gids,
                "run_ids": sorted({r for gid in gids for r in (groups[gid].get("run_ids") or [])}),
                "series": series,
                "candidate_forms": forms,
                "selection_rule": "per_seed_scatter is preferred: it shows every observation, so "
                                  "the sample size behind the claim cannot be hidden by a summary",
                "exploratory": any(s["kind"] == "exploratory" for s in series),
            })
        return ({"generated_at": time.time(), "items": items,
                 "no_figure_without_a_claim": True,
                 "refused": refusals}, refusals)

    def _resolve_groups(self, fig, claim_id, results) -> list[str]:
        """Which analysis groups this figure is entitled to draw."""
        wanted: list[str] = []
        for r in results.get("reported") or []:
            if str(claim_id) in [str(c) for c in (r.get("claim_ids") or [])]:
                if r.get("group_id"):
                    wanted.append(str(r["group_id"]))
        src = fig.get("data_source") or {}
        rids = {str(x) for x in (src.get("result_ids") or [])}
        metrics = {str(x) for x in (src.get("metrics") or [])}
        named = {str(x) for x in (src.get("group_ids") or [])}
        for g in results.get("groups") or []:
            gid = str(g.get("group_id"))
            run_ids = {str(x) for x in (g.get("run_ids") or [])}
            # A result id may name a run, a group, or the experiment/branch the group
            # belongs to — the manuscript plans figures per experiment, while the
            # analysis groups per (experiment, metric). Recognising the branch is an
            # identifier question, not a loosening: the numeric check below still has
            # to pass, so a figure bound this way can still be refused for disagreeing
            # with the numbers it claims to show.
            branch = str(g.get("branch") or "")
            # `experiment_id` is what a figure plan names; `branch` is that experiment's
            # arm. Both are accepted, because a figure plotted per experiment wants
            # every arm of it and a figure naming one arm wants only that one.
            exp = str(g.get("experiment_id") or "")
            if (gid in named or (rids and (rids & run_ids)) or (rids and branch in rids)
                    or (rids and exp and exp in rids)
                    or (metrics and str(g.get("metric")) in metrics)):
                wanted.append(gid)
        known = {str(g.get("group_id")) for g in results.get("groups") or []}
        resolved = sorted({w for w in wanted if w in known})

        # A branch usually carries several metrics, and a figure is drawn for ONE
        # claim, which names ONE metric. Leaving all of a branch's groups bound makes
        # the figure plot one metric and be verified against another — which is
        # exactly what the read-back check caught. The claim text is the authority on
        # which metric this figure is about, so narrow to it when it is unambiguous.
        by_metric: dict[str, list[str]] = {}
        for g in results.get("groups") or []:
            gid = str(g.get("group_id"))
            if gid in resolved:
                by_metric.setdefault(str(g.get("metric")), []).append(gid)
        if len(by_metric) > 1:
            text = f"{fig.get('caption') or ''} {fig.get('title') or ''} {claim_id or ''}".lower()
            named_metrics = [m for m in by_metric if m and m.lower() in text]
            if len(named_metrics) == 1:
                return sorted(by_metric[named_metrics[0]])
        return resolved

    def _asserted_numbers(self, fig) -> list[dict[str, Any]]:
        """Numbers the plan claims the figure will show, from caption or spec."""
        out = _quantities(str(fig.get("caption") or ""))
        for v in (fig.get("data_source") or {}).get("asserted_values") or []:
            if isinstance(v, (int, float)):
                text = repr(float(v))
                precision = len(text.split(".")[1]) if "." in text else 0
                out.append({"raw": text, "value": float(v), "precision": precision,
                            "percent": False})
        return out

    def _render(self, plt, item, form, style):
        """Draw the figure, then read the numbers back off the artists."""
        series = sorted(item["series"], key=lambda s: str(s["branch"]))
        width = float(style.get("width_in", 1.7 * max(3, len(series))))
        height = float(style.get("height_in", 3.4))
        fig, ax = plt.subplots(figsize=(width, height))
        elements: list[dict[str, Any]] = []
        drawn_points: dict[str, list[float]] = {}
        drawn_means: dict[str, float] = {}

        for i, s in enumerate(series):
            gid_stub = f"fig-{_slug(item['figure_id'])}-{_slug(str(s['branch']))}"
            if form == "per_seed_scatter":
                art = ax.scatter([i] * len(s["values"]), s["values"], s=30, zorder=3)
                art.set_gid(f"{gid_stub}-points")
                drawn_points[s["branch"]] = [float(y) for _x, y in art.get_offsets()]
                elements.append({"element_id": f"{gid_stub}-points", "role": "observations",
                                 "binds_to": {"artifact": "analysis_results",
                                              "group_id": s["group_id"], "statistic": "values"},
                                 "values": [float(v) for v in s["values"]]})
            if s["ci95"]:
                line, = ax.plot([i, i], s["ci95"], lw=3, alpha=0.35, zorder=2)
                line.set_gid(f"{gid_stub}-ci95")
                elements.append({"element_id": f"{gid_stub}-ci95", "role": "interval_95",
                                 "binds_to": {"artifact": "analysis_results",
                                              "group_id": s["group_id"], "statistic": "ci95"},
                                 "values": [float(v) for v in s["ci95"]]})
            if s["mean"] is not None:
                mean_art = ax.scatter([i], [s["mean"]], marker="_", s=460, zorder=4)
                mean_art.set_gid(f"{gid_stub}-mean")
                drawn_means[s["branch"]] = float(mean_art.get_offsets()[0][1])
                elements.append({"element_id": f"{gid_stub}-mean", "role": "mean",
                                 "binds_to": {"artifact": "analysis_results",
                                              "group_id": s["group_id"], "statistic": "mean"},
                                 "values": [float(s["mean"])]})
            label = ax.annotate(f"n={s['n']}", (i, s["mean"] if s["mean"] is not None else 0),
                                textcoords="offset points", xytext=(0, 10), ha="center",
                                fontsize=7)
            label.set_gid(f"{gid_stub}-n")
            elements.append({"element_id": f"{gid_stub}-n", "role": "sample_size",
                             "binds_to": {"artifact": "analysis_results",
                                          "group_id": s["group_id"], "statistic": "n"},
                             "values": [float(s["n"] or 0)]})

        ax.set_xticks(range(len(series)))
        ax.set_xticklabels([str(s["branch"]) for s in series], rotation=15, ha="right")
        ax.set_ylabel(str(item["metric"]))
        title = ax.set_title(f"{item['metric']} — {item['claim_id']}")
        title.set_gid(f"fig-{_slug(item['figure_id'])}-title")
        elements.append({"element_id": f"fig-{_slug(item['figure_id'])}-title", "role": "title",
                         "binds_to": {"artifact": "manuscript_spine", "claim_id": item["claim_id"]},
                         "values": []})
        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="svg")
        plt.close(fig)
        return buf.getvalue(), {"elements": elements, "points": drawn_points, "means": drawn_means}

    def _drift(self, drawn, item) -> str:
        """What the artists hold versus what the analysis says."""
        problems = []
        for s in item["series"]:
            expected = [float(v) for v in s["values"]]
            got = drawn["points"].get(s["branch"])
            if got is not None and (len(got) != len(expected) or
                                    not all(_close(a, b) for a, b in zip(sorted(got),
                                                                         sorted(expected)))):
                problems.append(f"{s['branch']}: plotted {got} for analysis values {expected}")
            mean = drawn["means"].get(s["branch"])
            if mean is not None and s["mean"] is not None and not _close(mean, float(s["mean"])):
                problems.append(f"{s['branch']}: plotted mean {mean} for analysis mean {s['mean']}")
        return "; ".join(problems)

    def _choose(self, variants, item) -> dict[str, Any]:
        by_form = {v["form"]: v for v in variants}
        return by_form.get("per_seed_scatter", variants[0])

    def _validate_svg(self, svg: bytes, elements) -> dict[str, Any]:
        """Editability and fidelity, checked on the file that will ship."""
        root = ET.fromstring(svg.decode("utf-8"))
        ids = {el.get("id") for el in root.iter() if el.get("id")}
        texts = [el for el in root.iter(f"{{{SVG_NS}}}text")]
        images = [el for el in root.iter(f"{{{SVG_NS}}}image")]
        missing = [e["element_id"] for e in elements if e["element_id"] not in ids]
        return {
            "is_svg": root.tag.endswith("svg"),
            "no_raster_embedded": not images,
            "text_is_real_text": bool(texts),
            "text_nodes": len(texts),
            "semantic_ids_present": not missing,
            "missing_ids": missing,
            "every_element_bound": all(e.get("binds_to") for e in elements),
        }

    def _reconstruct(self, ctx: Context):
        """Vector reconstruction of an externally supplied raster figure.

        A structured visual spec can be rebuilt as vector. A bare raster cannot —
        not honestly. Tracing it here would produce a file that looks editable and
        whose numbers came from a guess at pixel positions, so instead the raster
        is refused and the refusal is recorded.
        """
        raster = ctx.external("source_raster_figure")
        spec = ctx.external("visual_spec")
        if not raster and not spec:
            return ({"figures": [], "reconstructed": False,
                     "reason": "no external raster figure or visual spec was supplied"},
                    {"figures": [], "reconstructed": False}, None)
        if raster and not spec:
            reason = (f"a raster figure ({raster}) was supplied without a structured visual spec. "
                      f"Reconstructing it would mean reading its numbers off pixel positions and "
                      f"presenting the guess as an editable vector figure.")
            return ({"figures": [], "reconstructed": False, "refused": reason},
                    {"figures": [], "reconstructed": False, "refused": reason},
                    "OFFLINE-EQUIVALENT REFUSAL: " + reason)
        elements = spec.get("elements", []) if isinstance(spec, dict) else list(spec or [])
        mapping = [{"source_region": e.get("region"), "element_id": f"recon-{i:03d}",
                    "role": e.get("role"), "text": e.get("text"), "value": e.get("value")}
                   for i, e in enumerate(elements)]
        return ({"figures": [{"figure_id": "R-01", "elements": mapping,
                              "source": str(raster) if raster else "visual_spec"}],
                 "reconstructed": True},
                {"figures": [{"figure_id": "R-01", "mapping": mapping}], "reconstructed": True},
                None)

    def _captions(self, storyboard) -> str:
        L = ["# Figure captions (draft)", "",
             "Each caption names the claim its figure argues. A caption that states a number the "
             "analysis does not carry is refused at generation time, not corrected here.", ""]
        for it in storyboard["items"]:
            L += [f"**{it['figure_id']}** (claim `{it['claim_id']}`). {it['caption']} "
                  f"Metric: {it['metric']}; groups: {', '.join(it['group_ids'])}; "
                  f"{'exploratory' if it['exploratory'] else 'confirmatory'}.", ""]
        if not storyboard["items"]:
            L.append("_no figure survived the plan._")
        return "\n".join(L)

    def _validation_md(self, trace, element_map, results, venue) -> str:
        L = ["# SVG validation report", "",
             f"Target venue: **{venue}**. Renderer: vector only — no figure in this project is "
             f"rasterised at any point, and the selected file and the editable file are the same "
             f"bytes.", "",
             "| figure | claim | vector | text is text | no raster | ids | elements bound |",
             "|---|---|---|---|---|---|---|"]
        for t in trace:
            c = t["checks"]
            em = element_map[t["figure_id"]]
            L.append(f"| `{t['figure_id']}` | `{t['claim_id']}` | {c['is_svg']} | "
                     f"{c['text_is_real_text']} ({c['text_nodes']} nodes) | "
                     f"{c['no_raster_embedded']} | {c['semantic_ids_present']} | "
                     f"{c['every_element_bound']} ({len(em['elements'])}) |")
        if not trace:
            L.append("| — | — | — | — | — | — | — |")
        L += ["", "## Fidelity", "",
              f"Every plotted value was read back off the rendered artists and compared with "
              f"`analysis/analysis_results.json` (ledger digest "
              f"`{results.get('ledger_digest')}`). A mismatch is a hard refusal: no figure is "
              f"emitted, because a figure that disagrees with the analysis is believed before the "
              f"analysis is checked.", "",
              "## What this report does not claim", "",
              "That the figures are the right figures. It claims that the numbers in them are the "
              "numbers the audited analysis produced, and that the files remain editable."]
        return "\n".join(L)


# ==========================================================================
# deck-factory
# ==========================================================================
@register
class DeckFactory(Skill):
    """A native deck whose every number points at an artifact.

    Consolidates defense-ppt-storyline, paper-to-ppt-evidence-mapper and
    defense-ppt-generator. The mapper existed because a deck built from a paper
    loses the paper's evidence bindings; keeping it as a separate skill meant the
    generator could run without it. Here the binding is a precondition of writing
    the file: a slide carrying a number that resolves to no artifact stops the
    build.
    """

    name = "deck-factory"

    def execute(self, ctx: Context) -> SkillResult:
        gone = _missing(ctx, ("manuscript_spine", "draft_manifest", "manuscript_draft",
                              "experiment_ledger", "comparison_mode", "selected_figure",
                              "svg_element_map"))
        if gone:
            raise GateBlocked(
                "deck_inputs",
                f"deck-factory cannot build a defense deck: {gone} have not been produced. A "
                f"defense slide is where an unbacked number gets stated out loud to a room that "
                f"cannot check it.",
                "run manuscript-builder, experiment-runner and figure-factory first")
        prs_mod = self._pptx()
        spine = ctx.store.read(self.name, "manuscript_spine")
        ledger = ctx.store.read(self.name, "experiment_ledger", default=[])
        meta = ctx.store.read(self.name, "meta_analysis", default={})
        cm = ctx.store.read(self.name, "comparison_mode")
        figures = _dir_manifest(ctx, "selected_figure", default={"figures": []})
        element_map = ctx.store.read(self.name, "svg_element_map", default={"figures": {}})
        validation = ctx.store.read(self.name, "svg_validation_report", default="")
        editable = _dir_manifest(ctx, "editable_svg", default={"figures": []})
        audience = ctx.external("defense_audience_profile", "thesis committee")
        minutes = float(ctx.external("target_talk_duration_minutes", 20) or 20)
        template = ctx.external("pptx_template")
        warnings: list[str] = []

        if not ledger:
            raise GateBlocked(
                "deck_inputs",
                "the experiment ledger is empty, so no number on any slide could be traced to a "
                "run that happened.",
                "run the experiments before building a defense deck")

        sources = self._sources(ledger, meta)
        spec, dropped = self._spec(spine, ledger, meta, cm, figures, audience, minutes)
        if dropped:
            warnings.append(
                f"removed {len(dropped)} slide element(s) using claim language forbidden under "
                f"{cm['mode']}: {sorted({d['pattern'] for d in dropped})}. A defense slide is the "
                f"least correctable place for a comparison the project cannot make.")
        evidence, unbound = self._bind(spec, sources)
        if unbound:
            raise GateBlocked(
                "slide_evidence_binding",
                f"{len(unbound)} quantitative slide element(s) bind to no project artifact: " +
                "; ".join(f"{u['slide_id']}/{u['element_id']} states {u['numbers']}"
                          for u in unbound[:8]) +
                ". A number on a defense slide that resolves to no recorded run is a number the "
                "room cannot check and the speaker cannot defend.",
                "remove the number, or bind it to an experiment_ledger or meta_analysis value")

        deck_bytes, native = self._build(prs_mod, spec, evidence, template)
        ctx.store.write(self.name, "deck_spec", spec)
        ctx.store.write(self.name, "slide_evidence", evidence)
        ctx.store.write(self.name, "slide_outline", self._outline_md(spec, cm))
        ctx.store.write(self.name, "speaker_timing", self._timing_csv(spec))
        ctx.store.write(self.name, "defense_deck", deck_bytes)
        ctx.store.write(self.name, "speaker_notes", self._notes_md(spec, cm, audience))
        ctx.store.write(self.name, "deck_manifest", {
            "run_id": ctx.run_id, "generated_at": time.time(),
            "audience": audience, "target_minutes": minutes,
            "comparison_mode": cm["mode"],
            "slides": [{"slide_id": s["slide_id"], "title": s["title"], "kind": s["kind"],
                        "seconds": s["seconds"], "figure": s.get("figure")} for s in spec["slides"]],
            "native_objects": native,
            "full_slide_images": 0,
            "figures_shipped_as": "vector SVG alongside the deck",
            "figures": [{"figure_id": f["figure_id"], "file": f["file"],
                         "claim_id": f.get("claim_id"),
                         "element_map": bool((element_map.get("figures") or {}).get(f["figure_id"]))}
                        for f in figures.get("figures", [])],
            "editable_figures_present": bool(editable.get("figures")),
            "figure_validation_reviewed": bool(validation.strip()),
            "every_quantitative_element_bound": True,
            "removed_for_comparison_mode": dropped,
        })

        if not figures.get("figures"):
            warnings.append("no figure was selected upstream, so the deck carries none. An empty "
                            "figure set is reported rather than filled with decoration.")
        warnings.append(
            "figures ship beside the deck as editable SVG and are named on their slides rather "
            "than embedded: PowerPoint cannot hold this vector format, and rasterising them to "
            "embed them would destroy the only editable copy.")
        return SkillResult(
            self.name,
            produced=["deck_spec", "slide_evidence", "slide_outline", "speaker_timing",
                      "defense_deck", "speaker_notes", "deck_manifest"],
            warnings=warnings, next_state="WRITING",
            detail={"slides": len(spec["slides"]), "native_objects": native,
                    "bound_elements": len(evidence["elements"]),
                    "comparison_mode": cm["mode"]})

    # ------------------------------------------------------------------
    def _pptx(self):
        try:
            import pptx  # noqa: F401
            from pptx import Presentation  # noqa: F401
        except Exception as e:  # noqa: BLE001
            raise GateBlocked(
                "deck_renderer",
                f"python-pptx is unavailable ({e}). deck-factory will not emit a stand-in file "
                f"named defense.pptx: a deck of images, or an empty file with the right "
                f"extension, is worse than no deck because it is discovered in the room.",
                "pip install python-pptx") from e
        import pptx as mod
        return mod

    def _sources(self, ledger, meta) -> list[dict[str, Any]]:
        """Every number this project is entitled to put on a slide."""
        src: list[dict[str, Any]] = []
        for e in _ledger_index(ledger):
            src.append({"artifact": "experiment_ledger", "locator":
                        f"{e['experiment_id']}/{e['run_id']}#{e['metric']}", "value": e["value"]})
        for path, value in _walk_numbers(meta):
            src.append({"artifact": "meta_analysis", "locator": path, "value": float(value)})
        # Counts are facts about the ledger itself, and a defense deck states them
        # constantly ("three seeds", "nine runs"). They are bound to the ledger, not
        # waved through as "obviously fine".
        src.append({"artifact": "experiment_ledger", "locator": "count:runs",
                    "value": float(len(ledger))})
        branches: dict[str, int] = {}
        for r in ledger:
            b = str((r.get("provenance") or {}).get("branch", "main"))
            branches[b] = branches.get(b, 0) + 1
        for b, n in branches.items():
            src.append({"artifact": "experiment_ledger", "locator": f"count:runs[{b}]",
                        "value": float(n)})
        src.append({"artifact": "experiment_ledger", "locator": "count:branches",
                    "value": float(len(branches))})
        seeds = {str((r.get("provenance") or {}).get("seed")) for r in ledger}
        src.append({"artifact": "experiment_ledger", "locator": "count:seeds",
                    "value": float(len(seeds))})
        # Scalars the runs recorded that are not metrics — wall-clock, seed number.
        # A deck's run table shows these constantly. They ARE traceable to a run id,
        # so refusing them was the binder looking in the wrong place rather than the
        # number being unbacked. The rule is "every number resolves to a recorded
        # run", and these do.
        for r in ledger:
            rid = r.get("run_id", "?")
            eid = r.get("experiment_id", "?")
            for field in ("seconds", "seed"):
                v = r.get(field, (r.get("provenance") or {}).get(field))
                if isinstance(v, (int, float)) and not isinstance(v, bool):
                    src.append({"artifact": "experiment_ledger",
                                "locator": f"{eid}/{rid}#{field}", "value": float(v)})
        return src

    def _spec(self, spine, ledger, meta, cm, figures, audience, minutes):
        """The storyline, before any file is opened."""
        forbidden = [str(p) for p in cm.get("forbidden_claim_patterns") or []]
        dropped: list[dict[str, Any]] = []

        def keep(text: str, slide_id: str) -> bool:
            for pat in forbidden:
                if re.search(re.escape(pat), text, re.I):
                    dropped.append({"slide_id": slide_id, "pattern": pat, "text": text[:160]})
                    return False
            if cm["mode"] == "CM_NONE" and COMPARATIVE_RE.search(text):
                dropped.append({"slide_id": slide_id,
                                "pattern": COMPARATIVE_RE.search(text).group(0), "text": text[:160]})
                return False
            return True

        claims = [c for c in spine.get("claims") or [] if c.get("kind") != "disclosure"]
        slides: list[dict[str, Any]] = []

        slides.append({"slide_id": "S-01", "kind": "title",
                       "title": _norm(spine.get("thesis"))[:120] or "Defense",
                       "bullets": [f"Audience: {audience}",
                                   f"Comparison mode: {cm['mode']} "
                                   f"(from {cm.get('derived_from_level')})"],
                       "table": None, "figure": None})
        claim_bullets = [f"{c['claim_id']}: {_norm(c.get('statement'))[:150]}" for c in claims]
        slides.append({"slide_id": "S-02", "kind": "claims",
                       "title": "What this work claims",
                       "bullets": [b for b in claim_bullets if keep(b, "S-02")],
                       "table": None, "figure": None})

        header, rows = self._results_table(ledger)
        slides.append({"slide_id": "S-03", "kind": "results",
                       "title": "Measured results, every run",
                       "bullets": [f"{len(ledger)} recorded runs across "
                                   f"{len({str((r.get('provenance') or {}).get('branch', 'main')) for r in ledger})} branches"],
                       "table": {"header": header, "rows": rows}, "figure": None})

        for i, fig in enumerate(figures.get("figures") or []):
            sid = f"S-{4 + i:02d}"
            caption = _norm(fig.get("caption"))
            slides.append({"slide_id": sid, "kind": "figure",
                           "title": f"Evidence: {fig.get('claim_id')}",
                           "bullets": [b for b in [caption] if b and keep(b, sid)],
                           "table": None,
                           "figure": {"figure_id": fig.get("figure_id"), "file": fig.get("file"),
                                      "claim_id": fig.get("claim_id"),
                                      "ships_as": "figures/editable/" + str(fig.get("file")),
                                      "embedded_as_image": False}})

        disclosure = cm.get("disclosure_required") or {}
        n = len(slides) + 1
        if disclosure.get("required") and disclosure.get("text_template"):
            slides.append({"slide_id": f"S-{n:02d}", "kind": "disclosure",
                           "title": "Comparison scope",
                           "bullets": [_norm(disclosure["text_template"])],
                           "table": None, "figure": None})
            n += 1
        limits = [_norm(c.get("statement"))[:150] for c in claims if c.get("kind") == "negative"]
        slides.append({"slide_id": f"S-{n:02d}", "kind": "limitations",
                       "title": "What this does not show",
                       "bullets": limits or ["Negative and null results are reported in the "
                                             "manuscript, not omitted here."],
                       "table": None, "figure": None})
        n += 1
        slides.append({"slide_id": f"S-{n:02d}", "kind": "close",
                       "title": "Summary",
                       "bullets": [f"{c['claim_id']}: {_norm(c.get('statement'))[:120]}"
                                   for c in claims[:4]
                                   if keep(_norm(c.get('statement'))[:120], f"S-{n:02d}")],
                       "table": None, "figure": None})

        budget = minutes * 60.0
        weights = {"title": 0.5, "claims": 1.5, "results": 2.0, "figure": 1.5,
                   "disclosure": 1.0, "limitations": 1.0, "close": 1.0}
        total = sum(weights.get(s["kind"], 1.0) for s in slides)
        for s in slides:
            s["seconds"] = round(budget * weights.get(s["kind"], 1.0) / total, 1)
        return ({"generated_at": time.time(), "audience": audience,
                 "target_minutes": minutes, "comparison_mode": cm["mode"],
                 "slides": slides,
                 "policy": "native shapes only; no slide is an image of a slide"}, dropped)

    def _results_table(self, ledger):
        metrics = sorted({m for r in ledger for m, _v in _walk_numbers(r.get("metrics") or {})})
        header = ["run", "branch", "status"] + metrics
        rows = []
        for r in ledger:
            prov = r.get("provenance") or {}
            flat = dict(_walk_numbers(r.get("metrics") or {}))
            rows.append([str(r.get("run_id")), str(prov.get("branch", "main")),
                         str(r.get("status"))] +
                        [("" if flat.get(m) is None else f"{flat[m]:g}") for m in metrics])
        return header, rows

    def _bind(self, spec, sources):
        """Every number on a slide, resolved to the artifact it came from."""
        elements: list[dict[str, Any]] = []
        unbound: list[dict[str, Any]] = []
        for s in spec["slides"]:
            texts = [("title", s["title"])]
            texts += [(f"bullet[{i}]", b) for i, b in enumerate(s.get("bullets") or [])]
            if s.get("table"):
                for ri, row in enumerate(s["table"]["rows"]):
                    for ci, cell in enumerate(row):
                        texts.append((f"table[{ri}][{ci}]", str(cell)))
            if s.get("figure"):
                texts.append(("figure_caption", str(s["figure"].get("file"))))
            for element_id, text in texts:
                quantities = _quantities(str(text))
                bindings = []
                missing = []
                for q in quantities:
                    hit = _match_number(q, sources)
                    if hit is None:
                        missing.append(q["raw"])
                    else:
                        bindings.append({"number": q["raw"], "artifact": hit["artifact"],
                                         "locator": hit["locator"], "value": hit["value"]})
                rec = {"slide_id": s["slide_id"], "element_id": element_id,
                       "text": str(text)[:220], "quantitative": bool(quantities),
                       "bound_to": bindings}
                if s.get("figure") and element_id == "figure_caption":
                    rec["bound_to"].append({"artifact": "selected_figure",
                                            "locator": str(s["figure"].get("figure_id")),
                                            "value": None})
                elements.append(rec)
                if missing:
                    unbound.append({"slide_id": s["slide_id"], "element_id": element_id,
                                    "numbers": missing, "text": str(text)[:220]})
        return ({"generated_at": time.time(), "elements": elements,
                 "rule": "every quantitative slide element names the project artifact its number "
                         "came from; an unbound number stops the build",
                 "unbound": unbound,
                 "artifacts_referenced": sorted({b["artifact"] for e in elements
                                                 for b in e["bound_to"]})}, unbound)

    def _build(self, pptx_mod, spec, evidence, template):
        """Real PowerPoint objects: text frames, tables, shapes — never pictures."""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation(str(template)) if template else Presentation()
        by_slide: dict[str, list[dict[str, Any]]] = {}
        for e in evidence["elements"]:
            by_slide.setdefault(e["slide_id"], []).append(e)
        counts = {"text_frames": 0, "tables": 0, "pictures": 0, "notes": 0}

        # Index 5 is "Title Only" in the stock template; falling back to the first
        # layout keeps an unfamiliar external template usable instead of crashing.
        layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
        for s in spec["slides"]:
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text_frame.text = s["title"]
                counts["text_frames"] += 1
            else:
                box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9.0), Inches(1.0))
                box.text_frame.text = s["title"]
                counts["text_frames"] += 1
            top = Inches(1.7)
            if s.get("bullets"):
                body = slide.shapes.add_textbox(Inches(0.6), top, Inches(9.0), Inches(3.4))
                tf = body.text_frame
                tf.word_wrap = True
                for i, bullet in enumerate(s["bullets"]):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = f"- {bullet}"
                    para.font.size = Pt(16)
                counts["text_frames"] += 1
                top = Inches(1.7 + 0.35 * max(1, len(s["bullets"])))
            if s.get("table"):
                header, rows = s["table"]["header"], s["table"]["rows"]
                shape = slide.shapes.add_table(len(rows) + 1, len(header), Inches(0.6), top,
                                               Inches(9.0), Inches(0.4 * (len(rows) + 1)))
                table = shape.table
                for c, name in enumerate(header):
                    table.cell(0, c).text = str(name)
                for r, row in enumerate(rows, start=1):
                    for c, cell in enumerate(row):
                        table.cell(r, c).text = str(cell)
                counts["tables"] += 1
            if s.get("figure"):
                # The figure is named, not embedded. python-pptx cannot carry SVG, and
                # rasterising it to fit would destroy the only editable copy — which is
                # the entire point of producing vector figures upstream.
                box = slide.shapes.add_textbox(Inches(0.6), top, Inches(9.0), Inches(1.2))
                box.text_frame.text = (
                    f"Figure {s['figure']['figure_id']} ships as vector: "
                    f"{s['figure']['ships_as']} (not rasterised into this deck)")
                counts["text_frames"] += 1
            notes = slide.notes_slide.notes_text_frame
            notes.text = self._slide_notes(s, by_slide.get(s["slide_id"], []))
            counts["notes"] += 1
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue(), counts

    def _slide_notes(self, s, elements) -> str:
        bound = [f"{b['artifact']}:{b['locator']}" for e in elements for b in e["bound_to"]]
        lines = [f"{s['title']} ({s['kind']})", ""]
        lines += [f"- {b}" for b in s.get("bullets") or []]
        lines += ["", "Evidence on this slide: " + (", ".join(sorted(set(bound))) or "none "
                                                    "(no quantitative element)")]
        return "\n".join(lines)

    def _outline_md(self, spec, cm) -> str:
        L = ["# Defense outline", "",
             f"Comparison mode **{cm['mode']}** bounds every claim on every slide.", ""]
        for s in spec["slides"]:
            L += [f"## {s['slide_id']} — {s['title']} ({s['kind']}, {s['seconds']}s)", ""]
            L += [f"- {b}" for b in s.get("bullets") or []] or ["- (no text)"]
            L.append("")
        return "\n".join(L)

    def _timing_csv(self, spec) -> str:
        return _csv(["slide_id", "kind", "title", "seconds"],
                    [[s["slide_id"], s["kind"], s["title"], s["seconds"]]
                     for s in spec["slides"]])

    def _notes_md(self, spec, cm, audience) -> str:
        L = [f"# Speaker notes — {audience}", "",
             f"Total budget: {spec['target_minutes']} minutes. Comparison mode: "
             f"**{cm['mode']}**.", "",
             "Every number spoken from these slides is bound to a project artifact in "
             "`slides/slide_evidence.json`. If a question asks where a number came from, that "
             "file is the answer.", ""]
        for s in spec["slides"]:
            L += [f"## {s['slide_id']} — {s['title']} ({s['seconds']}s)", ""]
            L += [f"- {b}" for b in s.get("bullets") or []] or ["- (no text)"]
            if s["kind"] == "disclosure":
                L += ["", "State this slide verbatim. It is the sentence that makes every "
                          "comparison in the talk readable."]
            L.append("")
        return "\n".join(L)


# ==========================================================================
# release-gate
# ==========================================================================
#: Everything release-gate must be able to read before it may decide anything.
#: A release assembled from a partial view of the audits is a release that was
#: never audited.
REQUIRED_INPUTS = (
    "integrity_gate", "submission_blockers", "claim_audit", "citation_audit", "review_report",
    "stats_audit", "comparison_mode", "source_repro_report", "artifact_manifest",
    "provenance_log", "manuscript_draft", "bibliography", "findings", "finding_memory_graph",
    "defense_deck",
)

#: Artifacts that would physically ship. Each one needs a provenance record and a
#: digest before it may leave the project. The provenance log itself is read but is
#: not on this list: it grows with every subsequent write, so a digest recorded for
#: it is stale the moment the next artifact is written. It is the witness, not the
#: thing witnessed.
DELIVERABLES = (
    "manuscript_draft", "defense_deck", "bibliography", "findings", "finding_memory_graph",
    "claim_audit", "citation_audit", "integrity_gate", "submission_blockers", "review_report",
    "stats_audit", "source_repro_report", "comparison_mode", "artifact_manifest",
)

AI_DISCLOSURE = (
    "This work was produced with substantial assistance from an automated research system "
    "(ResearchForge). The system generated experiment code, executed runs, performed the "
    "statistical analysis, drafted the manuscript text and produced the figures and slides "
    "listed in this manifest. Every quantitative claim was checked mechanically against a "
    "recorded run. No part of this output has been verified by a human author unless an "
    "attestation appears in this manifest, and it is therefore released as an assisted draft "
    "rather than as a submission-ready manuscript."
)


@register
class ReleaseGate(Skill):
    """The last refusal.

    Was release-gate-exporter, and the rename is the point: exporting was the
    incidental part. Everything upstream of here has already decided whether the
    work is honest; this skill's only job is to make sure that decision is
    binding, that nothing ships without lineage, and that the reader is told the
    machine wrote it.
    """

    name = "release-gate"

    def execute(self, ctx: Context) -> SkillResult:
        gone = _missing(ctx, REQUIRED_INPUTS)
        if gone:
            raise GateBlocked(
                "release_inputs",
                f"release-gate cannot decide a release: {sorted(gone)} have not been produced. "
                f"A release assembled from a partial view of the audits was never audited.",
                "run the writing and review planes to completion before release")

        gate = ctx.store.read(self.name, "integrity_gate")
        blockers_md = ctx.store.read(self.name, "submission_blockers", default="")
        claim_audit = ctx.store.read(self.name, "claim_audit", default=[])
        citation_audit = ctx.store.read(self.name, "citation_audit", default="")
        review = ctx.store.read(self.name, "review_report", default={})
        stats = ctx.store.read(self.name, "stats_audit", default={})
        cm = ctx.store.read(self.name, "comparison_mode")
        repro = ctx.store.read(self.name, "source_repro_report", default={})
        manifest_in = ctx.store.read(self.name, "artifact_manifest", default={})
        provenance = ctx.store.read(self.name, "provenance_log", default=[])
        draft = ctx.store.read(self.name, "manuscript_draft", default="")
        findings = ctx.store.read(self.name, "findings", default=[])
        memory = ctx.store.read(self.name, "finding_memory_graph", default={})
        bib = ctx.store.read(self.name, "bibliography", default="")
        progress = ctx.store.read(self.name, "progress_state", default={})

        blockers: list[dict[str, Any]] = []
        blockers += self._carry_audits(gate, blockers_md, claim_audit, citation_audit,
                                       review, stats)
        blockers += self._enforce_comparison_mode(draft, cm, repro)
        blockers += self._refuse_synthetic(
            {"integrity_gate": gate, "review_report": review, "stats_audit": stats,
             "claim_audit": claim_audit, "findings": findings, "finding_memory_graph": memory,
             "artifact_manifest": manifest_in, "comparison_mode": cm,
             "source_repro_report": repro}, provenance)
        lineage, lineage_blockers = self._verify_provenance(ctx, manifest_in, provenance)
        blockers += lineage_blockers
        for i, b in enumerate(blockers, 1):
            b.setdefault("severity", "BLOCKER")
            b["blocker_id"] = f"R-{i:03d}"

        attested = self._attestations(gate, repro, review)
        released = not blockers
        status = "BLOCKED" if blockers else "ASSISTED_DRAFT"
        ai = {
            "statement": AI_DISCLOSURE,
            "system": "ResearchForge",
            "run_id": ctx.run_id,
            "mode": ctx.mode,
            "offline_provider_used": bool(ctx.offline),
            "machine_performed": sorted({str(e.get("skill")) for e in provenance
                                         if e.get("kind") == "artifact_write" and e.get("skill")}),
            "human_attestations": attested,
            "human_verified": bool(attested),
            # A system that writes manuscripts must not be the thing that decides they
            # are ready to submit. Major venues require the machine's role to be
            # declared, and "submission-ready" is a judgement only a named author can
            # make, so the default is an assisted draft and it stays that way here.
            "submission_ready": False,
            "why_not_submission_ready": (
                "release-gate never marks output submission-ready. The disclosure above must be "
                "carried into the venue's LLM-use declaration, and a named human author must "
                "verify the claims before submission."),
            "venue_policy_note": (
                "Most major venues now require disclosure of generative-model involvement, and "
                "several prohibit listing the system as an author. This manifest states the "
                "machine's role so that the declaration can be made accurately."),
        }
        deliverables, referenced = self._package(ctx, lineage, released)
        bundle = {
            "run_id": ctx.run_id, "decided_at": time.time(),
            "released": released, "release_status": status,
            "comparison_mode": cm["mode"],
            "deliverables": deliverables,
            "referenced_not_copied": referenced,
            "withheld": [] if released else [d["artifact_id"] for d in lineage["artifacts"]],
            "blockers": blockers,
            "ai_participation": ai,
            "progress_state_at_release": {"state": progress.get("state"),
                                          "run_id": progress.get("run_id")},
            "note": ("nothing was copied into release/ because the gate refused"
                     if not released else
                     "every file below has a provenance record and a verified digest"),
        }
        release_manifest = self._manifest(ctx, lineage, bundle, ai, cm, gate, released)
        self._validate_manifest(ctx, release_manifest)

        ctx.store.write(self.name, "release_bundle", bundle)
        ctx.store.write(self.name, "release_manifest", release_manifest)
        ctx.store.write(self.name, "release_report",
                        self._report_md(bundle, lineage, blockers, ai, cm, gate, repro, review,
                                        bib, findings))

        warnings: list[str] = []
        if blockers:
            warnings.append(
                f"RELEASE REFUSED: {len(blockers)} unresolved blocker(s). Nothing was packaged. "
                f"None of them is waivable by re-running this skill: "
                f"{sorted({b['kind'] for b in blockers})}")
        else:
            warnings.append(
                "released as an ASSISTED DRAFT, not as a submission-ready manuscript. The "
                "AI-participation disclosure in release_manifest.json must be carried into the "
                "venue's declaration, and a human author must verify the claims before "
                "submission.")
        return SkillResult(
            self.name, produced=["release_bundle", "release_manifest", "release_report"],
            warnings=warnings, next_state="RELEASE" if released else "REVIEWING",
            detail={"released": released, "release_status": status,
                    "blockers": len(blockers),
                    "blocker_kinds": sorted({b["kind"] for b in blockers}),
                    "artifacts_checked": len(lineage["artifacts"]),
                    "submission_ready": False})

    # ------------------------------------------------------------------
    def _carry_audits(self, gate, blockers_md, claim_audit, citation_audit, review, stats):
        """Upstream verdicts are carried, never re-litigated.

        release-gate does not re-audit a claim; it has no better view of the
        evidence than the auditor did. What it does is refuse to let an upstream
        BLOCK be forgotten between the audit and the export.
        """
        out: list[dict[str, Any]] = []
        for b in gate.get("blockers") or []:
            out.append({"kind": f"unresolved_{b.get('kind', 'blocker')}",
                        "source": f"integrity_gate:{b.get('blocker_id')}",
                        "detail": b.get("detail", ""),
                        "remediation": b.get("remediation", "resolve upstream and re-run")})
        if gate.get("verdict") == "BLOCK" and not (gate.get("blockers") or []):
            out.append({"kind": "integrity_gate_block", "source": "integrity_gate",
                        "detail": "the integrity gate returned BLOCK with no itemised blockers",
                        "remediation": "re-run claim-citation-auditor"})
        if gate.get("submission_permitted") is False and not out:
            out.append({"kind": "submission_not_permitted", "source": "integrity_gate",
                        "detail": "integrity_gate.submission_permitted is false",
                        "remediation": "resolve the audit findings"})
        md_blockers = re.findall(r"^##\s+`(B-\d+)`\s*(\S+)", blockers_md or "", re.M)
        if md_blockers and not (gate.get("blockers") or []):
            # The two artifacts come from the same skill. If they disagree, one of
            # them is stale, and a stale audit is not evidence of a clean one.
            out.append({"kind": "audit_disagreement", "source": "submission_blockers",
                        "detail": f"submission_blockers.md lists {len(md_blockers)} blocker(s) "
                                  f"that integrity_gate.json does not carry",
                        "remediation": "re-run claim-citation-auditor so both artifacts agree"})
        bad = [c for c in claim_audit
               if str(c.get("verdict")) in ("FABRICATED", "NOT_SUPPORTED", "SCOPE_MISMATCH")]
        if bad:
            out.append({"kind": "unsupported_claims_in_draft", "source": "claim_audit",
                        "detail": f"{len(bad)} audited claim(s) remain "
                                  f"{sorted({str(c['verdict']) for c in bad})}: "
                                  f"{[c.get('claim_id') or c.get('locator') for c in bad][:6]}",
                        "remediation": "remove or restate the claims, then re-run the audit"})
        for issue in (stats.get("findings") or stats.get("issues") or []):
            if isinstance(issue, dict) and str(issue.get("severity", "")).upper() in (
                    "BLOCKER", "CRITICAL", "HIGH"):
                out.append({"kind": "stats_audit_unresolved", "source": "stats_audit",
                            "detail": f"{issue.get('code') or issue.get('finding_id')}: "
                                      f"{issue.get('detail') or issue.get('message')}",
                            "remediation": issue.get("remediation", "resolve in integrity-auditor")})
        if (stats.get("evidence_lock") or {}).get("blocked"):
            out.append({"kind": "evidence_not_locked", "source": "stats_audit",
                        "detail": "the statistics audit refuses evidence lock for "
                                  f"{(stats['evidence_lock'].get('blocked_claims') or [])}",
                        "remediation": "resolve the blocking statistical findings"})
        if str(review.get("recommendation")) == "DO_NOT_SUBMIT":
            out.append({"kind": "review_do_not_submit", "source": "review_report",
                        "detail": str(review.get("recommendation_rationale") or
                                      "the review simulation recommends against submission"),
                        "remediation": "address the review findings"})
        if re.search(r"\bFABRICATED\b", citation_audit or ""):
            out.append({"kind": "fabricated_citation", "source": "citation_audit",
                        "detail": "the citation audit still reports a fabricated citation",
                        "remediation": "resolve or remove the citation"})
        return out

    def _enforce_comparison_mode(self, draft, cm, repro):
        """The degradation path must still bind at the last step.

        If the disclosure can be dropped between the audit and the export, or a
        comparative claim can survive CM_NONE by living in a file the auditor did
        not read, then the whole reproduction-level machinery is decorative.
        """
        out: list[dict[str, Any]] = []
        disclosure = cm.get("disclosure_required") or {}
        required_text = _norm(disclosure.get("text_template"))
        if disclosure.get("required"):
            if not required_text:
                out.append({"kind": "disclosure_undefined", "source": "comparison_mode",
                            "detail": f"{cm['mode']} requires a disclosure and none is defined",
                            "remediation": "re-run reproduction-fallback-planner"})
            elif required_text not in _norm(draft):
                out.append({"kind": "missing_disclosure", "source": "manuscript_draft",
                            "detail": f"comparison mode {cm['mode']} requires the disclosure "
                                      f"\"{required_text[:140]}\" and the released manuscript does "
                                      f"not contain it. Without it every comparison in the paper "
                                      f"reads as measured when it was not.",
                            "remediation": "insert the disclosure verbatim and re-run the audit"})
        if cm["mode"] == "CM_NONE":
            hits = sorted({m.group(0).lower() for m in COMPARATIVE_RE.finditer(draft or "")})
            if hits:
                out.append({"kind": "comparative_claim_under_CM_NONE",
                            "source": "manuscript_draft",
                            "detail": f"comparative language survives into the release under "
                                      f"CM_NONE ({hits[:6]}). At "
                                      f"{cm.get('derived_from_level')} no baseline was reproduced, "
                                      f"so there is nothing these comparisons were measured "
                                      f"against.",
                            "remediation": "remove the comparisons, or raise the reproduction "
                                           "level and re-derive the comparison mode"})
        if repro.get("level") == "RL0" and cm["mode"] != "CM_NONE":
            out.append({"kind": "mode_level_mismatch", "source": "comparison_mode",
                        "detail": f"the reproduction report is RL0 but the comparison mode is "
                                  f"{cm['mode']}; the weaker of the two governs",
                        "remediation": "re-run reproduction-fallback-planner"})
        return out

    def _refuse_synthetic(self, payloads, provenance):
        """Stub output must never be releasable.

        The offline provider exists so a run can proceed without a model; its
        output is a structural placeholder. Marked or not, it is not research, and
        this is the boundary where that stops being recoverable.
        """
        out: list[dict[str, Any]] = []
        hits: list[str] = []
        for label, payload in payloads.items():
            hits += _synthetic_hits(label, payload)
        for e in provenance:
            if e.get("kind") == "skill_end" and (e.get("detail") or {}).get("synthetic") is True:
                hits.append(f"provenance:{e.get('skill')}")
        if hits:
            out.append({"kind": "synthetic_artifact_in_release", "source": "provenance_log",
                        "detail": f"{len(hits)} artifact(s) or run(s) in this release chain are "
                                  f"marked synthetic: {sorted(set(hits))[:10]}. Output produced by "
                                  f"the offline stub is a placeholder, and a placeholder that "
                                  f"reaches a reader is a fabrication however carefully it was "
                                  f"labelled upstream.",
                        "remediation": "re-run the affected skills with a real model provider"})
        return out

    def _verify_provenance(self, ctx, manifest_in, provenance):
        """Every artifact in the release must have lineage and a digest.

        An artifact with no provenance record cannot be shown to have come from
        anything, so it cannot support a claim; an artifact whose digest does not
        match the one recorded when it was written has changed since it was
        audited. Both are blockers, not notes.
        """
        writes: dict[str, list[dict[str, Any]]] = {}
        reads_by_skill: dict[str, set[str]] = {}
        for e in provenance:
            if e.get("kind") == "artifact_write" and e.get("artifact_id"):
                writes.setdefault(str(e["artifact_id"]), []).append(e)
            if e.get("kind") == "artifact_read" and e.get("artifact_id"):
                reads_by_skill.setdefault(str(e.get("skill")), set()).add(str(e["artifact_id"]))

        rows: list[dict[str, Any]] = []
        blockers: list[dict[str, Any]] = []
        seen_paths: set[str] = set()
        for aid in DELIVERABLES:
            spec = ARTIFACTS[aid]
            path = ctx.store.path_for(aid)
            if spec.path.endswith("/"):
                path = path / "_manifest.json"
            events = writes.get(aid, [])
            recorded = events[-1].get("digest") if events else None
            actual = sha256_file(path) if path.exists() else None
            rel = str(path.relative_to(ctx.project)) if path.exists() else spec.path
            seen_paths.add(rel)
            row = {"artifact_id": aid, "path": rel, "kind": "deliverable",
                   "producer": spec.producer,
                   "recorded_sha256": recorded, "actual_sha256": actual,
                   "provenance_events": len(events),
                   "parents": sorted(reads_by_skill.get(spec.producer, set())),
                   "exists": bool(actual)}
            rows.append(row)
            if not events:
                blockers.append({
                    "kind": "artifact_without_provenance", "source": f"artifact:{aid}",
                    "detail": f"'{aid}' is in the release and the provenance log records no write "
                              f"of it. An artifact with no lineage cannot be shown to have come "
                              f"from anything, so it cannot support a claim.",
                    "remediation": f"re-run {spec.producer} so the write is recorded"})
            elif not recorded:
                blockers.append({
                    "kind": "artifact_without_digest", "source": f"artifact:{aid}",
                    "detail": f"'{aid}' has a provenance record with no digest, so what shipped "
                              f"cannot be shown to be what was audited.",
                    "remediation": f"re-run {spec.producer}"})
            elif actual is None:
                blockers.append({
                    "kind": "artifact_missing_on_disk", "source": f"artifact:{aid}",
                    "detail": f"'{aid}' is recorded in the provenance log but {rel} is not on "
                              f"disk.",
                    "remediation": f"re-run {spec.producer}"})
            elif actual != recorded:
                blockers.append({
                    "kind": "artifact_digest_drift", "source": f"artifact:{aid}",
                    "detail": f"'{aid}' has changed since it was written: recorded "
                              f"{str(recorded)[:16]}, on disk {str(actual)[:16]}. What would ship "
                              f"is not what was audited.",
                    "remediation": f"re-run {spec.producer} and re-run the audits over the result"})

        for entry in (manifest_in.get("artifacts") or []):
            aid = str(entry.get("artifact_id"))
            rel = str(entry.get("path"))
            declared = entry.get("sha256")
            p = ctx.project / rel
            actual = sha256_file(p) if p.is_file() else None
            events = writes.get(aid, [])
            rows.append({"artifact_id": aid, "path": rel,
                         "kind": str(entry.get("kind") or "referenced"),
                         "producer": (events[-1].get("skill") if events else None),
                         "recorded_sha256": declared, "actual_sha256": actual,
                         "provenance_events": len(events),
                         "parents": list(entry.get("parents") or []),
                         "exists": bool(actual)})
            if not events:
                blockers.append({
                    "kind": "artifact_without_provenance", "source": f"artifact_manifest:{aid}",
                    "detail": f"the experiment manifest lists '{aid}' ({rel}) and the provenance "
                              f"log records no write of it. An artifact with no lineage cannot "
                              f"support a claim.",
                    "remediation": "produce the artifact through the store, or remove it from "
                                   "artifact_manifest.json"})
            if not declared:
                blockers.append({
                    "kind": "artifact_without_digest", "source": f"artifact_manifest:{aid}",
                    "detail": f"the experiment manifest lists '{aid}' ({rel}) with no sha256.",
                    "remediation": "re-run experiment-runner so the digest is recorded"})
            elif actual is not None and actual != declared:
                blockers.append({
                    "kind": "artifact_digest_drift", "source": f"artifact_manifest:{aid}",
                    "detail": f"'{aid}' ({rel}) does not match the digest recorded for it.",
                    "remediation": "re-run experiment-runner"})
        return ({"artifacts": rows,
                 "complete": not blockers,
                 "rule": "an artifact in the release needs a provenance write event, a recorded "
                         "digest, and a file whose digest still matches"}, blockers)

    def _attestations(self, gate, repro, review) -> list[dict[str, Any]]:
        out = []
        for label, payload in (("integrity_gate", gate), ("source_repro_report", repro),
                               ("review_report", review)):
            if payload.get("human_reviewed") is True:
                out.append({"artifact": label, "attested_by": payload.get("reviewed_by"),
                            "at": payload.get("reviewed_at")})
        return out

    def _package(self, ctx, lineage, released):
        """Copy the deliverables in, but only if the gate let them through."""
        out_dir = ctx.store.path_for("release_bundle")
        out_dir.mkdir(parents=True, exist_ok=True)
        deliverables: list[dict[str, Any]] = []
        referenced: list[dict[str, Any]] = []
        for row in lineage["artifacts"]:
            if row["kind"] != "deliverable":
                referenced.append({"artifact_id": row["artifact_id"], "path": row["path"],
                                   "sha256": row["actual_sha256"],
                                   "note": "referenced by digest; not duplicated into the bundle"})
                continue
            if not released or not row["exists"]:
                continue
            src = ctx.project / row["path"]
            dst = out_dir / row["path"]
            dst.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(src, dst)
            deliverables.append({"artifact_id": row["artifact_id"],
                                 "path": str(dst.relative_to(ctx.project)),
                                 "sha256": row["actual_sha256"]})
        return deliverables, referenced

    def _manifest(self, ctx, lineage, bundle, ai, cm, gate, released) -> dict[str, Any]:
        artifacts = []
        for row in lineage["artifacts"]:
            artifacts.append({
                "artifact_id": row["artifact_id"],
                "path": row["path"],
                # An artifact with no digest is still listed. Omitting it would make
                # the manifest look complete, which is the failure this gate exists
                # to prevent; the empty digest is the finding.
                "sha256": row["actual_sha256"] or "",
                "kind": row["kind"],
                "parents": row["parents"],
                "producer": row["producer"],
                "provenance_events": row["provenance_events"],
                "provenance_complete": bool(row["provenance_events"] and row["actual_sha256"]
                                            and row["actual_sha256"] == row["recorded_sha256"]),
                "released": released and row["kind"] == "deliverable" and bool(row["actual_sha256"]),
            })
        return {
            "project_id": str(ctx.project.name),
            "artifacts": artifacts,
            "run_id": ctx.run_id,
            "generated_at": time.time(),
            "release_status": bundle["release_status"],
            "released": released,
            "submission_ready": False,
            "ai_participation": ai,
            "comparison_mode": {"mode": cm["mode"],
                                "derived_from_level": cm.get("derived_from_level"),
                                "disclosure_required": bool(
                                    (cm.get("disclosure_required") or {}).get("required"))},
            "integrity_gate": {"verdict": gate.get("verdict"),
                               "submission_permitted": gate.get("submission_permitted"),
                               "blockers": len(gate.get("blockers") or [])},
            "blockers": bundle["blockers"],
            "provenance_complete": lineage["complete"],
        }

    def _validate_manifest(self, ctx, manifest) -> None:
        """The release manifest is held to the same schema as any other artifact.

        Its contract entry declares no schema, so nothing would check it. A
        manifest of what shipped is the wrong artifact to leave unchecked.
        """
        from jsonschema import Draft202012Validator
        from jsonschema import ValidationError

        schema_path = Path(ctx.store.schemas_dir) / "ArtifactManifest.schema.json"
        try:
            Draft202012Validator(json.loads(schema_path.read_text())).validate(manifest)
        except ValidationError as e:
            raise SchemaViolation(
                f"release_manifest does not validate against ArtifactManifest: {e.message} "
                f"(at {'/'.join(str(x) for x in e.absolute_path) or '<root>'})") from e

    def _report_md(self, bundle, lineage, blockers, ai, cm, gate, repro, review, bib, findings):
        L = [f"# Release report — {bundle['release_status']}", "",
             f"Released: **{bundle['released']}**. Submission-ready: **false**.", "",
             f"- run: `{bundle['run_id']}`",
             f"- comparison mode: **{cm['mode']}** (from {cm.get('derived_from_level')}, "
             f"reproduction level {repro.get('level')})",
             f"- integrity gate: **{gate.get('verdict')}** "
             f"({len(gate.get('blockers') or [])} blocker(s))",
             f"- review recommendation: **{review.get('recommendation', 'n/a')}**",
             f"- findings carried: {len(findings)}",
             f"- bibliography entries: {bib.count('@')}", ""]
        if blockers:
            L += ["## Why this release was refused", "",
                  f"{len(blockers)} blocker(s). Nothing was copied into `release/`. None of these "
                  f"is waivable by re-running release-gate — the gate has no view of the evidence "
                  f"that the audits did not already have.", ""]
            for b in blockers:
                L += [f"### `{b['blocker_id']}` {b['kind']}", "",
                      f"- source: `{b['source']}`",
                      f"- finding: {b['detail']}",
                      f"- remediation: {b['remediation']}", ""]
        else:
            L += ["## Gate outcome", "",
                  "No blocker survived. This is not a statement that the work is correct. It is a "
                  "statement that every upstream audit passed, every released artifact has a "
                  "provenance record whose digest still matches the file on disk, nothing in the "
                  "chain is stub output, and the disclosure this comparison mode requires is "
                  "present.", ""]
        L += ["## Provenance completeness", "",
              "| artifact | kind | provenance events | digest matches | released |",
              "|---|---|---|---|---|"]
        for row in lineage["artifacts"]:
            matches = (row["actual_sha256"] is not None
                       and row["actual_sha256"] == row["recorded_sha256"])
            L.append(f"| `{row['artifact_id']}` | {row['kind']} | {row['provenance_events']} | "
                     f"{matches} | {bool(bundle['released']) and row['kind'] == 'deliverable'} |")
        L += ["", "## AI participation", "",
              f"> {ai['statement']}", "",
              f"- machine performed: {', '.join(ai['machine_performed']) or 'nothing recorded'}",
              f"- human attestations: {ai['human_attestations'] or 'none'}",
              f"- submission-ready: **false** — {ai['why_not_submission_ready']}",
              f"- venue policy: {ai['venue_policy_note']}", "",
              "## Contents", ""]
        L += [f"- `{d['artifact_id']}` -> `{d['path']}`" for d in bundle["deliverables"]] or \
             ["- nothing was packaged"]
        return "\n".join(L)
